from pptx import Presentation
import streamlit as st
from localApp import generate_text_from_model


def sum_pptx(pptx_file):
    st.warning('提交成功，教学PPT总结中，请稍后。。。。。。')

    prs = Presentation(pptx_file)

    txt = ""  # 初始化txt变量用于存储文本内容

    # 遍历每一页（幻灯片）
    for slide_number, slide in enumerate(prs.slides, start=1):
        txt += f"第 {slide_number} 页:\n"  # 添加页码信息

        # 遍历幻灯片中的所有形状
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt += shape.text + "\n"  # 添加形状的文本内容

    print(txt)

    prompt = """
    现有一个ppt，每页的内容如下，请对这个ppt进行总结，先说明每一页的内容，再总体概况，注意通常第一页为主标题和副标题
    """ + txt

    outputs = generate_text_from_model(prompt)

    st.success(outputs)
